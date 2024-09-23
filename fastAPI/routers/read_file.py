from fastapi import FastAPI, File, UploadFile, HTTPException
import fitz 
import pandas as pd
import json
import xml.etree.ElementTree as ET
import docx
import pptx
import os

def read_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def read_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()
    
def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

def read_json(file_content):
    return json.dumps(json.loads(file_content.decode('utf-8')), indent=4)

def read_xml(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    return ET.tostring(root, encoding='unicode', method='xml')

def read_docx(file_path):
    file_path = os.path.abspath(file_path)
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()
    
def read_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

if __name__ == "__main__":
    with open("./addresses.csv", "rb") as f:
        print(read_csv(f))
        print(f.filename)